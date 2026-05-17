"""
pptx_generator.py — Dynamic layout PPTX generator for Alert Normativo

Il template assets/_CLEAN.pptx contiene 6 slide identiche già pronte.
La logica di paginazione scorre le slide in ordine: quando current_y supera
BOTTOM_CONTENT_Y si passa alla slide successiva (reset current_y = TOP_CONTENT_Y).
Le slide vuote rimaste vengono lasciate intatte — le gestisce il responsabile.

Slide A4 portrait: 7,559,675 × 10,691,813 EMU

Content area:
  TOP_CONTENT_Y    = 2,065,553 EMU — fondo header + 100k margine
  BOTTOM_CONTENT_Y = 9,873,419 EMU — tetto footer

Colonne:
  Descrizione: left=254,373   width=5,325,951
  Data:        left=5,580,324 width=866,037
  Link:        left=6,446,361 width=866,037
"""

import os
import textwrap
from datetime import date
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from output.excel_logger import read_approved_news as _read_approved_news

# ── Struttura slide ───────────────────────────────────────────────────────────
SLIDE_W = 7559675
TOP_CONTENT_Y    = 2065553   # fondo "Principali aggiornamenti normativi" + 100k margine
BOTTOM_CONTENT_Y = 9873419   # tetto footer

# ── Colonne ───────────────────────────────────────────────────────────────────
DESC_LEFT  = 254373;  DESC_WIDTH  = 5325951
DATE_LEFT  = 5580324; DATE_WIDTH  = 866037
LINK_LEFT  = 6446361; LINK_WIDTH  = 866037

# ── Altezze ───────────────────────────────────────────────────────────────────
SECTION_HEADER_H = 303227   # 0.842 cm — intestazione sezione
SECTION_GAP      = 180000   # 0.500 cm — spazio tra sezioni diverse
ITEM_V_MARGIN    = 80000    # 0.222 cm — margine tra notizie consecutive

# ── Font e wrapping ───────────────────────────────────────────────────────────
FONT_NAME       = "Avenir Next LT Pro"
CHARS_PER_LINE  = 79
LINE_HEIGHT_EMU = int(Pt(14))   # 177,800 EMU — interlinea 1.4×
TEXT_V_PAD      = 36000         # padding verticale interno del textbox

# ── Icona sezione ─────────────────────────────────────────────────────────────
ICON_LEFT  = 34935
ICON_PAD   = 25614                             # (SECTION_HEADER_H - 252000) / 2
ICON_SIZE  = SECTION_HEADER_H - 2 * ICON_PAD  # ≈ 252,000 EMU
LABEL_LEFT  = ICON_LEFT + ICON_SIZE + 60000
LABEL_WIDTH = SLIDE_W - LABEL_LEFT - ICON_LEFT

CATEGORY_ORDER = ["BANKING", "INSURANCE", "CROSS FINANCE", "APPROFONDIMENTI"]
DATA_START_ROW  = 3


# ── Lettura Excel ─────────────────────────────────────────────────────────────

def read_approved_news(excel_path: str, edizione_numero: str) -> dict:
    """Wrapper — delegates to excel_logger.read_approved_news (DRY).

    Filters col H = SI and col J = edizione_numero.
    """
    return _read_approved_news(excel_path, edizione_numero=edizione_numero)


# ── Utilità slide ─────────────────────────────────────────────────────────────

def _update_edition(slide, numero: str, mese: str, anno: str) -> None:
    """Aggiorna il riquadro edizione (N. X / Mese Anno) nella slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "N. " in shape.text_frame.text and "\n" in shape.text_frame.text:
            runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
            if len(runs) >= 2:
                runs[0].text = f"N. {numero}"
                runs[1].text = f"{mese} {anno}"
            break


def _item_height(text: str) -> int:
    """Stima altezza EMU del textbox descrizione in base al wrapping del testo."""
    lines = textwrap.wrap(str(text), width=CHARS_PER_LINE) if text else []
    return max(1, len(lines)) * LINE_HEIGHT_EMU + TEXT_V_PAD


# ── Disegno elementi ──────────────────────────────────────────────────────────

def _draw_section_header(slide, section: str, icon_path: str, y: int) -> None:
    """Disegna icona PNG + label sezione in nero 14pt bold."""
    if os.path.exists(icon_path):
        slide.shapes.add_picture(
            icon_path,
            ICON_LEFT + ICON_PAD,
            y + ICON_PAD,
            ICON_SIZE,
            ICON_SIZE,
        )
    tb = slide.shapes.add_textbox(LABEL_LEFT, y, LABEL_WIDTH, SECTION_HEADER_H)
    tf = tb.text_frame
    tf.word_wrap = False
    tf._txBody.bodyPr.set("anchor", "ctr")
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = section
    run.font.name  = FONT_NAME
    run.font.size  = Pt(14)
    run.font.bold  = True
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def _draw_news_item(slide, item: dict, y: int, item_h: int) -> None:
    """Disegna descrizione, data e link cliccabile per una notizia."""
    BLACK = RGBColor(0x00, 0x00, 0x00)
    NAVY  = RGBColor(0x00, 0x20, 0x60)

    # Descrizione
    tb = slide.shapes.add_textbox(DESC_LEFT, y, DESC_WIDTH, item_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf._txBody.bodyPr.set("anchor", "t")
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = str(item.get("descrizione", ""))
    run.font.name  = FONT_NAME
    run.font.size  = Pt(10)
    run.font.bold  = False
    run.font.color.rgb = BLACK

    # Data
    tb = slide.shapes.add_textbox(DATE_LEFT, y, DATE_WIDTH, item_h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf._txBody.bodyPr.set("anchor", "t")
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = str(item.get("data", ""))
    run.font.name  = FONT_NAME
    run.font.size  = Pt(10)
    run.font.bold  = False
    run.font.color.rgb = BLACK

    # Link
    tb = slide.shapes.add_textbox(LINK_LEFT, y, LINK_WIDTH, item_h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf._txBody.bodyPr.set("anchor", "t")
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = "Link"
    run.font.name  = FONT_NAME
    run.font.size  = Pt(10)
    run.font.bold  = False
    run.font.color.rgb = NAVY
    if item.get("url"):
        run.hyperlink.address = item["url"]


# ── Entry point ───────────────────────────────────────────────────────────────

def generate_pptx(
    template_path: str,
    excel_path: str,
    output_path: str,
    numero: str,
    mese: str,
    anno: str,
    edizione_numero: str = "",
) -> None:
    """Genera la PPTX dal template (6 slide pre-esistenti) e dall'Excel revisionato.

    Paginazione: scorre le slide del template in ordine; quando current_y supera
    BOTTOM_CONTENT_Y passa alla slide successiva. Le slide vuote non vengono eliminate.

    edizione_numero: valore atteso in col J — se vuoto, usa `numero` come fallback.
    """
    grouped = read_approved_news(excel_path, edizione_numero=edizione_numero or numero)

    prs    = Presentation(template_path)
    slides = list(prs.slides)

    if not slides:
        raise ValueError(f"Il template {template_path!r} non contiene slide.")

    slide_idx = 0
    slide     = slides[slide_idx]
    _update_edition(slide, numero, mese, anno)
    current_y     = TOP_CONTENT_Y
    first_section = True

    def _next_slide():
        nonlocal slide_idx, slide, current_y
        slide_idx += 1
        if slide_idx >= len(slides):
            raise RuntimeError(
                f"Esaurite le slide del template ({len(slides)} disponibili). "
                "Aggiungi altre slide al template _CLEAN.pptx."
            )
        slide = slides[slide_idx]
        _update_edition(slide, numero, mese, anno)
        current_y = TOP_CONTENT_Y

    for section in CATEGORY_ORDER:
        items = grouped.get(section, [])
        if not items:
            continue

        icon_path = os.path.join(
            "assets",
            f"icon_{section.lower().replace(' ', '_')}.png",
        )

        if not first_section:
            current_y += SECTION_GAP
        first_section = False

        # Nuova slide se l'intestazione non entra
        if current_y + SECTION_HEADER_H > BOTTOM_CONTENT_Y:
            _next_slide()

        _draw_section_header(slide, section, icon_path, current_y)
        current_y += SECTION_HEADER_H

        for item in items:
            item_h = _item_height(item.get("descrizione", ""))

            # Nuova slide se la notizia non entra
            if current_y + item_h > BOTTOM_CONTENT_Y:
                _next_slide()

            _draw_news_item(slide, item, current_y, item_h)
            current_y += item_h + ITEM_V_MARGIN

    slides_used = slide_idx + 1
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    print(
        f"[OK] PPTX salvata: {output_path} "
        f"({slides_used} slide usate / {len(slides)} totali)"
    )


def get_output_path(numero: str, mese: str, anno: str) -> str:
    today = date.today().strftime("%Y%m%d")
    return f"output/ALERT_PPT/alert_normativo_N{numero}_{mese}{anno}_{today}.pptx"
